import copy
import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Pt


# 将多个单页PPTX组装成完整的PPTX
def compose_pptx_template(pptx_file_paths: list[str], output_path) -> Presentation:
    # 创建一个新的空白演示文稿
    merged_prs = Presentation()
    merged_prs.slide_width = Pt(960)   # 宽屏分辨率宽度
    merged_prs.slide_height = Pt(540)  # 宽屏分辨率高度

    # 由于默认会有1页空白幻灯片，先删除它
    if len(merged_prs.slides) == 1 and not merged_prs.slides[0].shapes:
        xml_slides = merged_prs.slides._sldIdLst
        slides = list(xml_slides)
        xml_slides.remove(slides[0])

    for pptx_file_path in pptx_file_paths:
        print(f"合并文件：{pptx_file_path}")
        src_prs = Presentation(pptx_file_path)

        for slide in src_prs.slides:
            # 使用空白布局添加一页新幻灯片
            new_slide = merged_prs.slides.add_slide(merged_prs.slide_layouts[6])

            for shape in slide.shapes:
                # 拷贝每个形状
                el = shape.element
                new_el = copy.deepcopy(el)
                new_slide.shapes._spTree.insert_element_before(new_el, 'p:extLst')

    # 保存合并后的新文件
    merged_prs.save(output_path)
    print(f"\n✅ 合并完成，保存为：{output_path}")

# 通过文本框替换文本（只需要处理普通PPT文件，不需要PPT模板提供占位符）
def fill_pptx_content(pptx_path, output_path, placeholders_to_replace):
    prs = Presentation(pptx_path)

    def replace_text_in_shapes(shapes, placeholders_to_replace_in_page):
        for shape in shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                # 是组合形状，递归处理
                replace_text_in_shapes(shape.shapes, placeholders_to_replace_in_page)
            elif shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        for placeholder_text, replacement_text in placeholders_to_replace_in_page.items():
                            if placeholder_text in run.text:
                                print(f"替换占位符: {placeholder_text} -> {replacement_text}")
                                run.text = run.text.replace(placeholder_text, replacement_text)

    for slide_idx, slide in enumerate(prs.slides, start=1):
        print(f"🔄 正在处理幻灯片 {slide_idx} ...")
        placeholders_to_replace_in_page = placeholders_to_replace[slide_idx - 1]
        replace_text_in_shapes(slide.shapes, placeholders_to_replace_in_page)

    prs.save(output_path)
    print(f"\n✅ 替换完成，已保存为：{output_path}")



##########################################
base_template_path = "ppt_templates"
file_name_list = [
    os.path.join(base_template_path, "real_ppt1.pptx"),
    # os.path.join(base_template_path, "t1.pptx"),
    # os.path.join(base_template_path, "t2.pptx"),
    # os.path.join(base_template_path, "t3.pptx"),
]
# 合并模板
compose_pptx_template(file_name_list, "1.pptx")

# 定义要替换的占位符文本和替换内容
placeholders_to_replace = [
    {
        "{{title}}": "项目开发过程",
        "{{text1}}": "需求分析",
        "{{text2}}": "方案设计",
        "{{text3}}": "开发实现",
        "{{text4}}": "测试验证",
        "{{text5}}": "上线部署",
        "{{text6}}": "维护优化",
    },
    {
        "{{title}}": "项目验收流程",
        "{{text1}}": "验收准备",
        "{{text2}}": "功能确认",
        "{{text3}}": "问题记录与修复",
        "{{text4}}": "文档与交付",
        "{{text5}}": "签署验收报告",
    },
    # {
    #     "{{title}}": "项目开发过程",
    #     "{{text1}}": "需求分析",
    #     "{{text2}}": "方案设计",
    #     "{{text3}}": "开发实现",
    #     "{{text4}}": "测试验证",
    #     "{{text5}}": "上线部署",
    #     "{{text6}}": "维护优化",
    # },
]

# 填入占位符文本
fill_pptx_content(os.path.join("1.pptx"), "2.pptx", placeholders_to_replace)