import copy
import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Pt

# 通过ppt母版文件对象和占位符对象渲染所有页ppt（母版占位符很难搞，后面估计会弃用这种方式）
def gen_content(prs, placeholders_to_replace):
    # 遍历每一张幻灯片
    for slide_number, slide in enumerate(prs.slides, start=1):
        print(f"幻灯片 {slide_number}:")
        placeholders_to_replace_in_page = placeholders_to_replace[slide_number - 1]
        # 遍历幻灯片中的所有占位符
        for placeholder in slide.placeholders:
            # 遍历预设的占位符和替换内容
            for placeholder_text, replacement_text in placeholders_to_replace_in_page.items():
                if placeholder_text in placeholder.text:
                    print(f"替换占位符文本: {placeholder_text} -> {replacement_text}")
                    placeholder.text = placeholder.text.replace(placeholder_text, replacement_text)
                    
        
            for shape in slide.shapes:
                # 检查是否是占位符并且是图片类型
                if shape.is_placeholder:
                    if (shape.placeholder_format.type == MSO_SHAPE_TYPE.PICTURE or shape.name.find('Picture') != -1) \
                        and not placeholder.text:
                        print(f"插入图片到占位符位置")
                        # 插入图片到占位符中
                        left = placeholder.left
                        top = placeholder.top
                        pic_path = "/Users/louisliu/Downloads/1.jpeg"  # 图片文件路径
                        pic = slide.shapes.add_picture(pic_path, left, top, width=placeholder.width, height=placeholder.height)
    return prs


# 将多个单页模板组装成完整模板
def gen_new_template(templates: list[str]) -> Presentation:
    # 创建一个新的 PowerPoint 文件
    new_template = Presentation()
    new_template.slide_width = Pt(960)   # 宽屏分辨率宽度
    new_template.slide_height = Pt(540)  # 宽屏分辨率高度
    for template in templates:
        tempPrs = Presentation(template)
        # 将模板中的所有幻灯片添加到新的演示文稿
        for slide in tempPrs.slides:
            slide_copy = new_template.slides.add_slide(slide.slide_layout)
            # 遍历原始幻灯片的形状
            for shape in slide.shapes:
                if shape.is_placeholder:  # 如果是占位符
                    # 获取当前占位符的类型（比如标题，占位符编号等）
                    placeholder_idx = shape.placeholder_format.idx
                    # 获取占位符的文本
                    if shape.has_text_frame:
                        text = shape.text  # 获取文本内容
                    else:
                        text = ''
                    # 尝试在新幻灯片中找到相应的占位符
                    try:
                        new_shape = slide_copy.shapes.placeholders[placeholder_idx]
                        # 复制占位符的大小和位置
                        new_shape.left = shape.left
                        new_shape.top = shape.top
                        new_shape.width = shape.width
                        new_shape.height = shape.height
                        # 复制文本内容
                        if new_shape.has_text_frame:
                            new_shape.text = text
                    except KeyError:
                        # 如果找不到相应的占位符，跳过或者做其他处理
                        print(f"Warning: No placeholder with idx {placeholder_idx} on the new slide.")
                        # 你可以选择继续或者用其他方法处理，比如用文本框替代占位符
                        pass
                # 处理图片占位符
                # elif shape.shape_type == 13:  # 图片类型
                #     # 复制图片占位符
                #     new_shape = slide_copy.shapes.add_picture(shape.image.filename, shape.left, shape.top, shape.width, shape.height)
    return new_template

# 通过文本框替换文本（只需要处理普通PPT文件，不需要PPT模板提供占位符）
def fill_pptx_content(pptx_path, output_path, placeholders_to_replace):
    prs = Presentation(pptx_path)

    def replace_text_in_shapes(shapes, placeholders_to_replace_in_page):
        for shape in shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:  # 是组合形状
                replace_text_in_shapes(shape.shapes)  # 递归处理子形状
            elif shape.has_text_frame:
                # 遍历预设的占位符和替换内容
                for placeholder_text, replacement_text in placeholders_to_replace_in_page.items():
                    if placeholder_text in shape.text:
                        print(f"替换占位符文本: {placeholder_text} -> {replacement_text}")
                        shape.text = replacement_text

    for slide_idx, slide in enumerate(prs.slides, start=1):
        print(f"幻灯片 {slide_idx}:")
        placeholders_to_replace_in_page = placeholders_to_replace[slide_idx - 1]
        replace_text_in_shapes(slide.shapes, placeholders_to_replace_in_page)

    # 保存新文件
    prs.save(output_path)
    print(f"\n✅ 已保存修改后的文件到：{output_path}")

##########################################
base_template_path = "ppt_templates"
file_name_list = [
    os.path.join(base_template_path, "all.pptx"),
    # os.path.join(base_template_path, "t1.pptx"),
    # os.path.join(base_template_path, "t2.pptx"),
    # os.path.join(base_template_path, "t3.pptx"),
]
# 合并模板
t = gen_new_template(file_name_list)
t.save("1.pptx")

# 定义要替换的占位符文本和替换内容
placeholders_to_replace = [
    {
        "{{title}}": "项目开发过程",
        "{{text1}}": "需求分析",
        "{{text2}}": "方案设计",
        "{{text3}}": "开发实现",
        "{{text4}}": "测试验证",
        "{{text5}}": "上线部署",
        "{{text6}}": "维护优化",
    },
    {
        "{{title}}": "项目验收流程",
        "{{text1}}": "验收准备",
        "{{text2}}": "功能确认",
        "{{text3}}": "问题记录与修复",
        "{{text4}}": "文档与交付",
        "{{text5}}": "签署验收报告",
    },
    # {
    #     "{{title}}": "项目开发过程",
    #     "{{text1}}": "需求分析",
    #     "{{text2}}": "方案设计",
    #     "{{text3}}": "开发实现",
    #     "{{text4}}": "测试验证",
    #     "{{text5}}": "上线部署",
    #     "{{text6}}": "维护优化",
    # },
]

# 渲染合成过后的模板
prs = Presentation("1.pptx")
prs = gen_content(prs, placeholders_to_replace)
prs.save("2.pptx")

