from pptx import Presentation

def print_placeholders_info(pptx_file):
    # 打开PPTX文件
    prs = Presentation(pptx_file)
    
    # 遍历所有幻灯片
    for slide_idx, slide in enumerate(prs.slides):
        print(f"幻灯片 {slide_idx + 1}:")
        
        # 遍历当前幻灯片中的所有形状
        for shape_idx, shape in enumerate(slide.shapes):
            if shape.is_placeholder:
                # 获取占位符的类型和名称
                placeholder_type = shape.placeholder_format.type
                placeholder_name = f"占位符 {placeholder_type}"
                
                # 打印占位符的详细信息
                print(f"  - 占位符 {shape_idx + 1}: 类型={placeholder_name}, 坐标=({shape.left}, {shape.top}), 大小=({shape.width}, {shape.height}), 文本='{shape.text if shape.has_text_frame else '无'}'")
                
# 调用函数，指定PPTX文件路径
pptx_file = "ppt_templates/t1.pptx"  # 替换为你的PPTX文件路径
print_placeholders_info(pptx_file)