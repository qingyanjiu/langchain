from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR

# 创建PPT
prs = Presentation()
slide_layout = prs.slide_layouts[5]  # 空白幻灯片
slide = prs.slides.add_slide(slide_layout)

# 添加标题
title_shape = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
title_frame = title_shape.text_frame
title_frame.text = "🌸 春游行程安排 🌸"
title_frame.paragraphs[0].font.size = Pt(36)
title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
title_frame.paragraphs[0].font.bold = True
title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 153, 0)  # 绿色

# 行程安排数据
schedule = [
    ("08:00", "集合出发"),
    ("09:00", "乘车前往景区"),
    ("10:30", "参观风景区"),
    ("12:00", "午餐时间"),
    ("14:00", "自由活动"),
    ("16:00", "返程回校")
]

# PPT 页面尺寸
slide_width = prs.slide_width
slide_height = prs.slide_height

# 形状大小
box_width = Inches(2.5)
box_height = Inches(1)
spacing = Inches(0.5)  # 形状之间的间隔

columns = 2  # 每行最多放两个
num_rows = (len(schedule) + columns - 1) // columns  # 计算总行数

# **动态计算起始位置，使整个流程居中**
total_height = num_rows * box_height + (num_rows - 1) * spacing
start_top = (slide_height - total_height) / 2  # 居中对齐

shapes = []
connectors = []  # 存放连接线

for i, (time, event) in enumerate(schedule):
    row = i // columns  # 当前行号
    col = i % columns  # 当前列号

    # **计算每行的中心位置**
    row_width = columns * box_width + (columns - 1) * spacing
    start_left = (slide_width - row_width) / 2  # 使该行居中

    # **蛇形排列：偶数行从左到右，奇数行从右到左**
    if row % 2 == 0:
        current_left = start_left + col * (box_width + spacing)  # 从左到右
    else:
        current_left = start_left + (columns - 1 - col) * (box_width + spacing)  # 从右到左

    current_top = start_top + row * (box_height + spacing)

    # **存储形状位置**
    shapes.append((current_left, current_top))

# **先绘制连接线（保证在底层）**
for i in range(len(shapes) - 1):
    start_x = shapes[i][0] + box_width / 2
    start_y = shapes[i][1] + box_height / 2
    end_x = shapes[i + 1][0] + box_width / 2
    end_y = shapes[i + 1][1] + box_height / 2

    if (i + 1) % columns == 0:
        # **换行时使用垂直线**
        connector = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT, start_x, start_y + box_height / 2, end_x, end_y - box_height / 2
        )
    else:
        # **同一行时**
        if (i // columns) % 2 == 0:
            # 偶数行：从左到右
            connector = slide.shapes.add_connector(
                MSO_CONNECTOR.STRAIGHT, start_x + box_width / 2, start_y, end_x - box_width / 2, end_y
            )
        else:
            # 奇数行：从右到左（缩短连接线）
            connector = slide.shapes.add_connector(
                MSO_CONNECTOR.STRAIGHT, start_x - box_width / 2 + spacing / 2, start_y, end_x + box_width / 2 - spacing / 2, end_y
            )

    connector.line.color.rgb = RGBColor(255, 165, 0)  # 橙色
    connector.line.width = Pt(3)  # 线宽
    connectors.append(connector)

# **调整连接线的 Z-Order，使其置于底层**
for connector in connectors:
    slide.shapes._spTree.remove(connector._element)  # 先移除
    slide.shapes._spTree.insert(2, connector._element)  # 重新插入到较低层

# **再绘制矩形（确保在连接线上方）**
for i, (time, event) in enumerate(schedule):
    current_left, current_top = shapes[i]

    # **添加矩形（行程节点）**
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, current_left, current_top, box_width, box_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(135, 206, 250)  # 天蓝色
    shape.line.color.rgb = RGBColor(0, 102, 204)  # 深蓝色边框

    # **在矩形内添加时间和事件**
    text_frame = shape.text_frame
    text_frame.text = f"{time}\n{event}"
    text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    text_frame.paragraphs[0].font.size = Pt(14)
    text_frame.paragraphs[0].font.bold = True

# **保存PPT**
prs.save("春游PPT_蛇形_连接线底层.pptx")
print("PPT 生成完成！")
