from pptx import Presentation

# 加载PPT文件
prs = Presentation("春游PPT_蛇形_连接线底层.pptx")

# 定义要替换的占位符文本和替换内容
placeholders_to_replace = {
    "{{title}}": "春游活动",
    "{{text}}": "这次春游活动将在春天举行，大家可以享受自然的美丽和放松的时光。",
}


# 遍历每一张幻灯片
for slide_number, slide in enumerate(prs.slides, start=1):
    print(f"幻灯片 {slide_number}:")
    
    # 遍历幻灯片中的所有占位符
    for placeholder in slide.placeholders:
        # 遍历预设的占位符和替换内容
        for placeholder_text, replacement_text in placeholders_to_replace.items():
            if placeholder_text in placeholder.text:
                print(f"替换占位符文本: {placeholder_text} -> {replacement_text}")
                placeholder.text = placeholder.text.replace(placeholder_text, replacement_text)
                
        # 检查占位符是否为图片占位符 (通常为 index 13，但需要根据实际验证)
        if not placeholder.text:  # 没有text就是图片占位符
            print(f"插入图片到占位符位置")
            
            # 插入图片到占位符中
            left = placeholder.left
            top = placeholder.top
            pic_path = "/Users/louisliu/Downloads/1.jpeg"  # 图片文件路径
            pic = slide.shapes.add_picture(pic_path, left, top, width=placeholder.width, height=placeholder.height)

# 保存修改后的PPT
prs.save("修改后的PPT.pptx")